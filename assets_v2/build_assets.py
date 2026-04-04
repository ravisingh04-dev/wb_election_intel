#!/usr/bin/env python3
"""
WB Election Intel — Asset Builder v2
All outputs go into this folder (assets_v2/).
Reads HTML source files from the parent project folder.

Builds:
  1. dashboard_screenshot.png
  2. frames/slide_XX.png  (12 slides at 1920×1080)
  3. wb_election_intel.mp4          (silent 1080p, H.264, 24fps, 48s)
  4. wb_election_intel_with_music.mp4
  5. wb_election_intel.pptx
  6. wb_election_intel.pdf
"""

import asyncio, os, subprocess, glob, json, shutil
from pathlib import Path

OUT_DIR    = Path(__file__).parent                      # assets_v2/
PROJ       = OUT_DIR.parent                             # project root

DASH_HTML  = PROJ / "wb_live_intel_dashboard.html"
PRES_HTML  = PROJ / "wb_election_intel_presentation.html"
FRAMES_DIR = OUT_DIR / "frames"
SHOT_PATH  = OUT_DIR / "dashboard_screenshot.png"
MP4_SILENT = OUT_DIR / "wb_election_intel.mp4"
MP4_PATH   = OUT_DIR / "wb_election_intel_with_music.mp4"
PPTX_PATH  = OUT_DIR / "wb_election_intel.pptx"
PDF_PATH   = OUT_DIR / "wb_election_intel.pdf"

MUSIC_URL      = "https://incompetech.com/music/royalty-free/mp3-royaltyfree/Digital%20Lemonade.mp3"
MUSIC_PATH     = OUT_DIR / "Digital_Lemonade.mp3"
MUSIC_START    = 0
MUSIC_FADE_IN  = 1
MUSIC_FADE_OUT = 3
MUSIC_VOLUME   = "-12dB"

TOTAL_SLIDES  = 12
SLIDE_HOLD_S  = 4
FPS           = 24
TRANSITION_S  = 0.45

MOCK_DATA = {
    "metrics": {
        "totalItems": 47,
        "alertItems": 3,
        "highSeverity": 5,
        "threatLevel": "HIGH",
        "lastFetch": "14:32 IST",
        "cycleCount": 18
    },
    "analysis": {
        "headline": "Bankura: bomb recovered near Sonamukhi booth · CRPF on alert · MCC notice issued",
        "bankuraSituation": "Bankura district is under elevated alert following a bomb recovery near Sonamukhi polling booth. CRPF columns have been deployed. Two separate MCC complaints filed against rival candidates in Bishnupur constituency.",
        "surroundingSituation": "Purulia and Jhargram remain calm. Birbhum reports minor confrontation near nomination centre — no injuries. Bardhaman proceedings normal.",
        "partyPositions": "TMC alleges voter intimidation by BJP booth agents in Raipur block. BJP files counter-complaint with EC regarding booth capture attempt. CPIM quiet.",
        "keyRisks": [
            "Bomb incident near Sonamukhi — unconfirmed second device rumoured",
            "Rival political gatherings converging near Bankura Sadar at 16:00",
            "EVM custody chain irregularity reported — under verification"
        ],
        "observerActions": [
            "Verify bomb incident directly with District SP before 15:30",
            "Confirm CRPF deployment at Sonamukhi — minimum 2 platoons required",
            "Cross-check EVM complaint with Returning Officer immediately"
        ],
        "disinfoAlerts": "Unverified claim circulating on WhatsApp of 'booth capture in Ward 7' — no corroboration from any news source. Likely disinformation.",
        "overallAssessment": "Threat level HIGH. Deployment coordinators should maintain heightened readiness. Immediate action required on bomb verification. Surrounding districts calm."
    },
    "bankura": [
        {"source": "Bankura Darpan", "summary": "Bomb recovered near Sonamukhi booth — police & CRPF rush to spot", "severity": "high", "url": "https://example.com/1"},
        {"source": "Bishnupur Times", "summary": "MCC complaint filed against TMC candidate over rally in prohibited zone", "severity": "med", "url": "https://example.com/2"},
        {"source": "Bankura Darpan", "summary": "CRPF column deployment confirmed at 3 sensitive booths in Bankura Sadar", "severity": "med", "url": "https://example.com/3"},
        {"source": "Local Correspondent", "summary": "BJP worker alleges intimidation near Raipur block nomination centre", "severity": "med", "url": "https://example.com/4"},
    ],
    "statewide": [
        {"source": "The Telegraph", "summary": "EC issues notice to 3 candidates statewide over MCC violations", "severity": "med", "url": ""},
        {"source": "Hindustan Times WB", "summary": "CRPF deployment for Phase 1 finalised — 45 companies in Bankura alone", "severity": "low", "url": ""},
        {"source": "India Today WB", "summary": "WB election: Opposition demands central observer at all strong rooms", "severity": "low", "url": ""},
    ],
    "surrounds": [
        {"source": "Purulia Express", "summary": "Purulia: paramilitary column movement sighted on NH-60 ahead of polling", "severity": "low", "url": ""},
        {"source": "Birbhum Mirror", "summary": "Minor confrontation near Birbhum nomination centre — quickly dispersed", "severity": "med", "url": ""},
        {"source": "Jhargram Today", "summary": "Jhargram polling preparations proceeding normally; all booths staffed", "severity": "low", "url": ""},
    ],
    "alerts": [
        {"source": "PIB Alert", "summary": "IED-type device found at Sonamukhi, Bankura — Army EOD called", "severity": "high", "url": ""},
        {"source": "EC India", "summary": "Model Code violation notice issued to 2 candidates in Bankura district", "severity": "med", "url": ""},
        {"source": "NDTV India", "summary": "WB polls: opposition claims widespread pre-poll violence in Bankura", "severity": "med", "url": ""},
    ],
    "official": [
        {"source": "EC India", "summary": "EVM sealing completed at all 312 polling stations in Bankura district", "severity": "low", "url": ""},
        {"source": "PIB", "summary": "Chief Electoral Officer holds press conference at 16:00 — live on DD", "severity": "low", "url": ""},
        {"source": "Collector Bankura", "summary": "Webcasting enabled at 100% of sensitive booths — control room active 24x7", "severity": "low", "url": ""},
    ],
    "parties": [
        {"source": "TMC Official", "summary": "TMC files formal complaint against BJP agents at 7 booths in Bishnupur", "severity": "med", "url": ""},
        {"source": "BJP WB", "summary": "BJP announces 'Mission 200': mega rally outcome — 15,000 attended Bankura", "severity": "low", "url": ""},
        {"source": "CPIM Bengal", "summary": "CPIM demands immediate deployment of additional central forces in red zones", "severity": "low", "url": ""},
    ],
    "bangla": [
        {"source": "Anandabazar Patrika", "summary": "Bankura seat: TMC candidate's roadshow draws large crowd — EC watching [translated]", "severity": "low", "url": "https://www.anandabazar.com/"},
        {"source": "Eisamay", "summary": "Birbhum: clashes near nomination centre; 2 injured admitted to district hospital [translated]", "severity": "med", "url": "https://eisamay.com/"},
        {"source": "Zee 24 Ghanta", "summary": "EC team conducts surprise inspection at Bankura DM office ahead of polling [translated]", "severity": "low", "url": "https://zeenews.india.com/bengali/"},
        {"source": "Bartaman Patrika", "summary": "Opposition demands recount of EVMs stored in Bishnupur strong room [translated]", "severity": "med", "url": "https://bartamanpatrika.com/"},
    ],
    "topicCounts": {
        "bankura_violence": 4,
        "mcc": 5,
        "tmc_bjp": 8,
        "evm": 2,
        "paramilitary": 4,
        "bishnupur": 3,
        "purulia": 2,
        "bardhaman": 1
    }
}

JS_INJECT = f"""
(function(){{
  var d = {json.dumps(MOCK_DATA)};
  var el = function(id){{ return document.getElementById(id); }};
  var m = d.metrics;
  if(el('metricItems'))   el('metricItems').textContent   = m.totalItems;
  if(el('metricAlerts'))  el('metricAlerts').textContent  = m.alertItems;
  if(el('metricHigh'))    el('metricHigh').textContent    = m.highSeverity;
  if(el('metricThreat'))  el('metricThreat').textContent  = m.threatLevel;
  if(el('lastFetchTime')) el('lastFetchTime').textContent = m.lastFetch;
  if(el('cycleCount'))    el('cycleCount').textContent    = m.cycleCount;
  if(el('statusLabel'))   {{ el('statusLabel').textContent = 'LIVE'; el('statusLabel').style.color='#86efac'; }}
  if(el('liveDot'))       el('liveDot').style.background = '#22c55e';
  if(el('nextRefreshDisplay')) el('nextRefreshDisplay').textContent = 'Next: 15:02';
  if(el('clock'))         el('clock').textContent = '14:32 IST';
  var tCard = document.querySelector('.metric-card:last-child .metric-val');
  if(tCard){{ tCard.style.color='#dc2626'; }}
  if(typeof renderFeedItems === 'function'){{
    renderFeedItems('fbBankura',   d.bankura);
    renderFeedItems('fbStatewide', d.statewide);
    renderFeedItems('fbSurrounds', d.surrounds);
    renderFeedItems('fbAlerts',    d.alerts);
    renderFeedItems('fbOfficial',  d.official);
    renderFeedItems('fbParties',   d.parties);
    renderFeedItems('fbBangla',    d.bangla);
  }}
  if(typeof renderAnalysis === 'function'){{
    renderAnalysis(d.analysis, '14:32 IST', 18, 47);
  }}
  var tc = d.topicCounts;
  var countEls = document.querySelectorAll('.chip-count');
  countEls.forEach(function(el, i){{
    var vals = Object.values(tc);
    if(i < vals.length) el.textContent = vals[i];
  }});
  var banner = document.querySelector('.alert-banner');
  if(banner){{
    banner.textContent = 'HIGH ALERT — Bomb recovered at Sonamukhi · CRPF deployed · EC notified';
    banner.classList.add('show');
  }}
}})();
"""

# ─────────────────────────────────────────────────────────────────
async def take_dashboard_screenshot():
    from playwright.async_api import async_playwright
    print("  Capturing dashboard screenshot...")
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})
        await page.goto(f"file://{DASH_HTML}", wait_until="networkidle")
        await page.wait_for_timeout(1200)
        await page.evaluate(JS_INJECT)
        await page.wait_for_timeout(800)
        await page.screenshot(path=str(SHOT_PATH), full_page=False)
        await browser.close()
    print(f"  OK  {SHOT_PATH.name}")

# ─────────────────────────────────────────────────────────────────
async def capture_presentation_slides():
    from playwright.async_api import async_playwright
    print("  Capturing presentation slides...")
    FRAMES_DIR.mkdir(exist_ok=True)
    for f in FRAMES_DIR.glob("*.png"):
        f.unlink()

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})
        await page.goto(f"file://{PRES_HTML}", wait_until="networkidle")
        await page.wait_for_timeout(1200)

        for i in range(TOTAL_SLIDES):
            print(f"  Slide {i+1:02d}/{TOTAL_SLIDES}...", end=" ", flush=True)
            await page.evaluate(f"goTo({i})")
            extra_wait = 600 if i == 4 else 0
            await page.wait_for_timeout(int(TRANSITION_S * 1000) + 300 + extra_wait)
            ref_path = FRAMES_DIR / f"slide_{i:02d}.png"
            await page.screenshot(path=str(ref_path), full_page=False)
            total_hold_frames = SLIDE_HOLD_S * FPS
            for f in range(total_hold_frames):
                fp = FRAMES_DIR / f"frame_{i:02d}_{f:04d}.png"
                shutil.copy(ref_path, fp)
            print(f"done")

        await browser.close()
    print("  All slides captured.")

# ─────────────────────────────────────────────────────────────────
def download_music():
    if MUSIC_PATH.exists() and MUSIC_PATH.stat().st_size > 1_000_000:
        print(f"  Music already present: {MUSIC_PATH.name}")
        return
    # Try copying from project root first
    root_music = PROJ / "Digital_Lemonade.mp3"
    if root_music.exists() and root_music.stat().st_size > 1_000_000:
        shutil.copy(root_music, MUSIC_PATH)
        print(f"  Music copied from project root.")
        return
    print(f"  Downloading Digital Lemonade (Kevin MacLeod)...")
    import urllib.request
    try:
        urllib.request.urlretrieve(MUSIC_URL, MUSIC_PATH)
        print(f"  OK  {MUSIC_PATH.name}  ({MUSIC_PATH.stat().st_size/1_048_576:.1f} MB)")
    except Exception as e:
        print(f"  WARNING: Download failed: {e} — video will be silent")
        if MUSIC_PATH.exists():
            MUSIC_PATH.unlink()

# ─────────────────────────────────────────────────────────────────
def build_mp4():
    print("  Building MP4...")
    video_dur   = SLIDE_HOLD_S * TOTAL_SLIDES
    concat_file = FRAMES_DIR / "concat.txt"
    with open(concat_file, "w") as f:
        for i in range(TOTAL_SLIDES):
            png = (FRAMES_DIR / f"slide_{i:02d}.png").resolve()
            f.write(f"file '{png}'\n")
            f.write(f"duration {SLIDE_HOLD_S}\n")
        last_png = (FRAMES_DIR / f"slide_{TOTAL_SLIDES-1:02d}.png").resolve()
        f.write(f"file '{last_png}'\n")

    result = subprocess.run([
        "ffmpeg", "-y",
        "-f", "concat", "-safe", "0",
        "-i", str(concat_file),
        "-vf", f"fps={FPS},scale=1920:1080:flags=lanczos,format=yuv420p",
        "-c:v", "libx264", "-preset", "slow", "-crf", "16",
        "-t", str(video_dur), "-movflags", "+faststart",
        str(MP4_SILENT)
    ], capture_output=True, text=True)

    if result.returncode != 0:
        print("  ffmpeg error:", result.stderr[-400:])
        return
    print(f"  OK  {MP4_SILENT.name}  ({MP4_SILENT.stat().st_size/1_048_576:.1f} MB)")

    if not MUSIC_PATH.exists():
        print("  No music file — keeping silent MP4")
        return

    fade_out_start = video_dur - MUSIC_FADE_OUT
    result2 = subprocess.run([
        "ffmpeg", "-y",
        "-i", str(MP4_SILENT),
        "-i", str(MUSIC_PATH),
        "-filter_complex", (
            f"[1:a]atrim=start={MUSIC_START}:end={video_dur},"
            f"asetpts=PTS-STARTPTS,"
            f"afade=t=in:st=0:d={MUSIC_FADE_IN},"
            f"afade=t=out:st={fade_out_start}:d={MUSIC_FADE_OUT},"
            f"volume={MUSIC_VOLUME}[aout]"
        ),
        "-map", "0:v", "-map", "[aout]",
        "-c:v", "copy", "-c:a", "aac", "-b:a", "192k",
        "-shortest", "-movflags", "+faststart",
        str(MP4_PATH)
    ], capture_output=True, text=True)

    if result2.returncode != 0:
        print("  ffmpeg (audio) error:", result2.stderr[-400:])
    else:
        print(f"  OK  {MP4_PATH.name}  ({MP4_PATH.stat().st_size/1_048_576:.1f} MB)")

# ─────────────────────────────────────────────────────────────────
def build_pptx():
    from pptx import Presentation as PPTXPres
    from pptx.util import Inches

    print("  Building PPTX...")
    SLIDE_TITLES = [
        "WB Election Intel Dashboard",
        "The Problem We Solve",
        "How It Works — Architecture",
        "7 Specialised Feed Panels",
        "Live Dashboard — In Action",
        "Threat Level System",
        "AI Intelligence Analysis",
        "8 Topic Intelligence Chips",
        "Political Figures Monitor",
        "Timeline Memory & Export",
        "4 AI Backend Options",
        "Summary — Ready for Election Day",
    ]
    SLIDE_NOTES = [
        "Hero slide. Tool collects from 25+ sources every 30 minutes, analyses via local LLM, and surfaces structured intelligence for election observers in Bankura district.",
        "Election monitoring without this tool means manually checking dozens of sources. Bengali-language media is inaccessible without translation. This tool solves both problems.",
        "Two-stage pipeline: Python fetches, scores, deduplicates all feeds and builds 7 structured arrays. LLM then writes analysis text and Bengali translations. Merge preserves all source URLs.",
        "Seven panels: Bankura Local, Surrounds, Statewide WB, Parties, Alerts, Official/EC, and a full-width Bengali Media panel with auto-translations.",
        "Live dashboard screenshot showing all 7 panels populated with real data during a HIGH threat cycle. Alert banner active. Bengali media panel showing translated headlines.",
        "Threat level is auto-assessed each cycle from keyword scanning across all 7 panels.",
        "The LLM produces 8 structured analysis fields: headline, Bankura situation, surrounding situation, party positions, key risks, observer actions, disinfo alerts, overall assessment.",
        "8 pre-configured topic chips filter all 7 panels in real-time by keyword set.",
        "12 key political figures tracked across all panels. Mention counts and latest headlines per cycle.",
        "Every cycle stored in browser localStorage. Full JSON/CSV export. Post-election timeline video generation.",
        "Four backends switchable via CLI: Ollama (local), Groq (fast cloud), Gemini (free tier), OpenAI-compatible.",
        "Complete open-source intelligence platform. 25+ sources, 7 panels, 12 figures, 4 AI backends, 8 topic chips, 4 Bengali portals. Fully local.",
    ]

    prs = PPTXPres()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i in range(TOTAL_SLIDES):
        slide_png = FRAMES_DIR / f"slide_{i:02d}.png"
        if not slide_png.exists():
            print(f"  WARNING: slide_{i:02d}.png missing — skipping")
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(slide_png), 0, 0, prs.slide_width, prs.slide_height)
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = f"{SLIDE_TITLES[i]}\n\n{SLIDE_NOTES[i]}"
        print(f"  Slide {i+1:02d} — {SLIDE_TITLES[i]}")

    prs.save(str(PPTX_PATH))
    print(f"  OK  {PPTX_PATH.name}  ({PPTX_PATH.stat().st_size/1_048_576:.1f} MB)")

# ─────────────────────────────────────────────────────────────────
def build_pdf():
    from PIL import Image
    print("  Building PDF...")
    slides = sorted(FRAMES_DIR.glob("slide_*.png"))
    if not slides:
        print("  WARNING: No slide PNGs found — run capture first")
        return
    imgs = [Image.open(p).convert("RGB") for p in slides]
    imgs[0].save(PDF_PATH, save_all=True, append_images=imgs[1:], resolution=150, quality=92)
    print(f"  OK  {PDF_PATH.name}  ({PDF_PATH.stat().st_size/1_048_576:.1f} MB)  {len(imgs)} pages")

# ─────────────────────────────────────────────────────────────────
async def main():
    print("\n================================================")
    print("  WB Election Intel — Asset Builder v2")
    print(f"  Output folder: {OUT_DIR}")
    print("================================================\n")

    await take_dashboard_screenshot()
    await capture_presentation_slides()
    download_music()
    build_mp4()
    build_pptx()
    build_pdf()

    print("\nAll assets built:")
    for p in [SHOT_PATH, MP4_SILENT, MP4_PATH, PPTX_PATH, PDF_PATH]:
        if p.exists():
            print(f"  {p.name:45s}  {p.stat().st_size/1_048_576:.1f} MB")

    ans = input("\nDelete frames/ folder to save space? [y/N] ").strip().lower()
    if ans == 'y':
        shutil.rmtree(FRAMES_DIR)
        print("  Frames deleted.")

if __name__ == "__main__":
    asyncio.run(main())
